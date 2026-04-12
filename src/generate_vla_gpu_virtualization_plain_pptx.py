#!/usr/bin/env python3

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/root/autodl-tmp/VLAServing")
RESULTS = ROOT / "results"
DOCS = ROOT / "docs"
OUT = DOCS / "VLA_Workload_GPU_Virtualization_Plain_20260413.pptx"

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)


def load_json(name: str):
    with (RESULTS / name).open() as f:
        return json.load(f)


VS = load_json("unified_chunked_vla_vs_baselines_20260412.json")
EFF = load_json("unified_chunked_vla_effectiveness_20260412.json")
PI2550 = load_json("pi05_vla_serving_autoh25_50_phase_shift_20260413.json")["pi05_autoh25_50"]
POLICY = load_json("vla_gpu_virtualization_policy_20260412.json")
ART = load_json("public_gpu_serving_artifacts_20260411.json")


def fmt(x: float, digits: int = 2) -> str:
    return f"{x:.{digits}f}"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    box = add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub = add_textbox(slide, Inches(0.62), Inches(1.05), Inches(12.0), Inches(0.5))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT


def add_lines(slide, left, top, width, height, lines: list[str], size: int = 20):
    box = add_textbox(slide, left, top, width, height)
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
    return box


def add_section_block(slide, left, top, width, title: str, lines: list[str], title_size: int = 18, body_size: int = 16):
    t = add_textbox(slide, left, top, width, Inches(0.35))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = BLACK
    b = add_textbox(slide, left, top + Inches(0.42), width, Inches(5.4))
    tf = b.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(body_size)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)


def add_page_num(slide, idx: int, total: int):
    box = add_textbox(slide, Inches(12.2), Inches(7.0), Inches(0.6), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.text = f"{idx}/{total}"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(10)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.RIGHT


def add_plain_rule(slide, top: float):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(12.0), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = []

    pi_fixed4 = PI2550["strict_horizon"]["fixed4_best"]["metrics"]
    pi_adm = PI2550["strict_horizon"]["admission"]["summary"]
    pi_old_adm = POLICY["pi05"]["strict_horizon"]["admission"]["summary"]
    pi_latest = VS["pi05"]["latest_method"]
    pi_full = VS["pi05"]["generic_full_resident_upper_bound"]
    pi_best = VS["pi05"]["best_conventional_baseline"]
    pi_gpulet = VS["pi05"]["gpulet_like"]

    gr_latest4 = VS["gr00t_n1d6"]["latest_method_4_robot"]
    gr_latest16 = VS["gr00t_n1d6"]["latest_method_16_robot_scaleout"]
    gr_full = VS["gr00t_n1d6"]["generic_full_resident_upper_bound"]
    gr_best = VS["gr00t_n1d6"]["best_conventional_baseline"]
    gr_gpulet = VS["gr00t_n1d6"]["gpulet_like"]
    gr_fair = VS["gr00t_n1d6"]["fair_admission"]
    gr_phase = EFF["gr00t_phase_lock_effectiveness"]["cases"]
    gr_mps = EFF["gr00t_batch_mps_effectiveness"]["cases"]

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "VLA Workload 的 GPU 虚拟化", "纯文本版：方法、实验、对照、结论")
    add_plain_rule(slide, Inches(1.35))
    add_lines(
        slide,
        Inches(0.8),
        Inches(1.8),
        Inches(11.8),
        Inches(4.6),
        [
            "目标：解释为什么已有 GPU 虚拟化 / 资源划分方案不足，",
            "以及为什么我们的 VLA-workload-aware 方法对 Pi05 和 GR00T N1.6 都有效。",
            "",
            "内容来源：",
            "- VLA_WORKLOAD_GPU_VGPU_ABSTRACTION.md",
            "- VLA_GPU_VIRTUALIZATION_POLICY_20260412.md",
            "- UNIFIED_CHUNKED_VLA_SERVING_SYSTEM.md",
            "- VLA_METHOD_REPRO_20260411.md",
            "- PUBLIC_ARTIFACT_REPRO_20260411.md",
        ],
        size=22,
    )
    slides.append(slide)

    # 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "VLA Workload 与传统 Serving 的差别")
    add_plain_rule(slide, Inches(1.2))
    add_section_block(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.9),
        "VLA workload 的关键特征",
        [
            "- 下一次请求大体可预测，不是完全外生随机到达",
            "- 真正的 deadline 是旧 action chunk 耗尽前必须返回新 chunk",
            "- 模型状态本身就是调度对象，不只是算力",
            "- horizon 可以被控制语义约束，也可以被系统主动利用",
            "- Pi05 当前使用 {25, 50} 控制语义，legal replan window 为 [25, 50]",
        ],
    )
    add_section_block(
        slide,
        Inches(6.75),
        Inches(1.45),
        Inches(5.8),
        "传统 serving / vGPU 默认假设",
        [
            "- 请求到达后再排队，再调度",
            "- 主要资源是 compute time / compute space / memory size",
            "- 目标偏向吞吐、平均延迟、SLO miss rate",
            "- 很少把未来请求时间线、模型状态驻留、带宽、合法重规划窗口当成一等资源",
        ],
    )
    slides.append(slide)

    # 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "我们的抽象：VLA-vGPU")
    add_plain_rule(slide, Inches(1.2))
    add_lines(
        slide,
        Inches(0.9),
        Inches(1.55),
        Inches(11.7),
        Inches(0.5),
        ["VLA-vGPU = <T, S, M, X, W, B, F>"],
        size=26,
    )
    add_section_block(
        slide,
        Inches(0.85),
        Inches(2.1),
        Inches(5.8),
        "7 个维度",
        [
            "- T: future compute slots",
            "- S: spatial compute share",
            "- M: resident model state",
            "- X: H2D / activation / apply bandwidth budget",
            "- W: legal replan window",
            "- B: batch affinity",
            "- F: fairness / admission weight",
        ],
        body_size=17,
    )
    add_section_block(
        slide,
        Inches(6.8),
        Inches(2.1),
        Inches(5.7),
        "核心思想",
        [
            "- 不是给每个机器人一块静态 GPU 分区",
            "- 而是给每个机器人一个随时间续租的 Robot Lease",
            "- 这个 lease 同时承诺 compute、model state、bandwidth 和 control window",
            "- admission 的对象是 robot lease，不是一次 request",
        ],
        body_size=17,
    )
    slides.append(slide)

    # 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "统一系统：一个控制面，两类后端")
    add_plain_rule(slide, Inches(1.2))
    add_section_block(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.8),
        "共享控制面",
        [
            "- Timeline Predictor",
            "- Admission Controller",
            "- Memory State Manager",
            "- Compute Scheduler",
            "- Phase Controller",
            "- 数据面统一成 compute / prefetch / decode-apply 三条流",
        ],
    )
    add_section_block(
        slide,
        Inches(6.75),
        Inches(1.45),
        Inches(5.8),
        "两类后端",
        [
            "- Pi05: ResidentPrefetch",
            "  关键在 predictive residency / predictive prefetch / window-aware serving",
            "- GR00T N1.6: SharedPrefixPhaseBatch",
            "  关键在 shared-prefix residency / phase lock / same-model batching / fair admission",
            "- 两者共享同一个 VLA-vGPU 抽象，只是 T/M/X/W/B 的重心不同",
        ],
    )
    slides.append(slide)

    # 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Pi05 如何套进这套方法")
    add_plain_rule(slide, Inches(1.2))
    add_section_block(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.8),
        "Pi05 运行时设计",
        [
            "- three-shell：30Hz 常驻、20Hz 常驻、两个 10Hz 共享",
            "- 10Hz 模型按 next-use 做 predictive prefetch",
            "- 请求窗口使用 {25, 50} horizon 语义",
            "- legal replan window 是 [25, 50]",
            "- 在 Pi05 上，真正关键的是 M + X + W，而不是 batching",
        ],
    )
    add_section_block(
        slide,
        Inches(6.75),
        Inches(1.45),
        Inches(5.8),
        "为什么有效",
        [
            "- 窗口变宽后，reservation 和 prefetch 更容易安排",
            "- 不需要等请求到了再 reactively 切模型",
            "- request-to-result 延迟基本不变",
            "- admission 容量显著提升",
        ],
    )
    slides.append(slide)

    # 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Pi05 实验结果")
    add_plain_rule(slide, Inches(1.2))
    add_section_block(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.8),
        "固定 4 机器人",
        [
            f"- request-to-result p95 = {fmt(pi_fixed4['mean_request_to_result_p95_ms'])} ms",
            f"- hard miss = {pi_fixed4['hard_miss_count']}",
            f"- reply-over = {pi_fixed4['reply_over_chunk_actions_count']}",
            f"- fleet score = {fmt(pi_fixed4['mean_fleet_score'], 3)}",
            f"- min robot score = {fmt(pi_fixed4['mean_min_robot_score'], 3)}",
        ],
    )
    add_section_block(
        slide,
        Inches(6.75),
        Inches(1.45),
        Inches(5.8),
        "admission + 对照",
        [
            f"- 当前 25/50 设定：mean admitted = {fmt(pi_adm['mean_admitted_total'])}",
            f"- 旧设定：mean admitted = {fmt(pi_old_adm['mean_admitted_total'])}",
            f"- admission 提升：{fmt(pi_adm['mean_admitted_total'] - pi_old_adm['mean_admitted_total'])}",
            f"- fleet score: {fmt(pi_old_adm['mean_fleet_score'], 4)} -> {fmt(pi_adm['mean_fleet_score'], 4)}",
            f"- min robot score: {fmt(pi_old_adm['mean_min_robot_score'], 4)} -> {fmt(pi_adm['mean_min_robot_score'], 4)}",
            f"- miss_autohorizon_ratio: {fmt(pi_old_adm['mean_miss_autohorizon_ratio'], 4)} -> {fmt(pi_adm['mean_miss_autohorizon_ratio'], 4)}",
        ],
    )
    slides.append(slide)

    # 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "GR00T N1.6 如何套进这套方法")
    add_plain_rule(slide, Inches(1.2))
    add_section_block(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.8),
        "GR00T 运行时设计",
        [
            "- 保留 shared-prefix residency，这个方法已验证有效",
            "- 相同微调模型的请求做 phase lock",
            "- 相同模型请求优先组成 batch",
            "- admission 使用 quota-fair，抑制 batch bias",
            "- 在 GR00T 上，关键是 T + M(shared prefix) + B + F",
        ],
    )
    add_section_block(
        slide,
        Inches(6.75),
        Inches(1.45),
        Inches(5.8),
        "为什么有效",
        [
            "- GR00T 的 same-model batching 很强",
            "- strict horizon 下会出现 reply-over",
            "- phase lock 能把请求对齐成稳定 batch",
            "- shared prefix 让多微调模型可以长期常驻",
            "- MPS 不是默认收益点，实验表明它在当前 runtime 上没有额外增益",
        ],
    )
    slides.append(slide)

    # 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "GR00T N1.6 实验结果")
    add_plain_rule(slide, Inches(1.2))
    add_section_block(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.8),
        "shared-prefix + phase-lock",
        [
            f"- 4 robots: p95 = {fmt(gr_latest4['mean_request_to_result_p95_ms'])} ms, stable = {gr_latest4['stable_under_100ms']}",
            f"- 8 robots: strict reply-over {gr_phase[0]['strict_horizon']['reply_over_chunk_actions_count']} -> phase-lock {gr_phase[0]['phase_lock_batch']['reply_over_chunk_actions_count']}",
            f"- 8 robots: batch {fmt(gr_phase[0]['strict_horizon']['mean_batch_size'],2)} -> {fmt(gr_phase[0]['phase_lock_batch']['mean_batch_size'],2)}",
            f"- 16 robots: p95 = {fmt(gr_latest16['mean_request_to_result_p95_ms'])} ms, stable = {gr_latest16['stable_under_100ms']}",
            f"- 16 robots: strict reply-over {gr_phase[1]['strict_horizon']['reply_over_chunk_actions_count']} -> phase-lock {gr_phase[1]['phase_lock_batch']['reply_over_chunk_actions_count']}",
            f"- 16 robots: batch {fmt(gr_phase[1]['strict_horizon']['mean_batch_size'],2)} -> {fmt(gr_phase[1]['phase_lock_batch']['mean_batch_size'],2)}",
        ],
        body_size=15,
    )
    add_section_block(
        slide,
        Inches(6.75),
        Inches(1.45),
        Inches(5.8),
        "fair admission + MPS 结论",
        [
            f"- accept-rate gap: {fmt(gr_fair['baseline_greedy_accept_rate_gap'],4)} -> {fmt(gr_fair['quota_fair_accept_rate_gap'],4)}",
            f"- final-count gap: {fmt(gr_fair['baseline_greedy_final_count_gap'],0)} -> {fmt(gr_fair['quota_fair_final_count_gap'],0)}",
            f"- 16 robots: batch-only p95 = {fmt(gr_mps[0]['batch_only']['mean_request_to_result_p95_ms'])} ms",
            f"- 16 robots: batch+MPS p95 = {fmt(gr_mps[0]['batch_plus_mps']['mean_request_to_result_p95_ms'])} ms",
            "- 结论：fair admission 有效；MPS 在当前 runtime 上没有额外收益",
        ],
        body_size=15,
    )
    slides.append(slide)

    # 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "为什么已有 GPU 虚拟化 / 资源划分方案不足")
    add_plain_rule(slide, Inches(1.2))
    add_section_block(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.8),
        "方法层面的根本问题",
        [
            "- GPUlet: offline integer-rate schedule，不建模 model state / prefetch / legal replan window",
            "- Clockwork: 只做 compute reservation，冷模型仍在关键路径上",
            "- REEF: preemption 不解决模型状态移动",
            "- Paella: reactive scheduling 只在模型已常驻时有效",
            "- USHER: 空间划分忽略未来请求可预测性和模型状态",
            "- DistServe: 对单 GPU 上的短 VLA 前向增加额外阶段队列",
        ],
        body_size=15,
    )
    add_section_block(
        slide,
        Inches(6.75),
        Inches(1.45),
        Inches(5.8),
        "实验和 artifact 证据",
        [
            f"- Pi05: ours {fmt(pi_latest['service_e2e_p95_ms'])} ms; full resident {fmt(pi_full['mean_latency_p95_ms'])} ms; best conventional {fmt(pi_best['mean_latency_p95_ms'])} ms",
            f"- GR00T: ours {fmt(gr_latest16['mean_request_to_result_p95_ms'])} ms at 16 robots; full resident {fmt(gr_full['mean_latency_p95_ms'])} ms; best conventional {fmt(gr_best['mean_latency_p95_ms'])} ms",
            f"- Pi05 GPUlet temporal feasible = {pi_gpulet['temporal_only']['feasible_under_100ms']}",
            f"- GR00T GPUlet temporal feasible = {gr_gpulet['temporal_only']['feasible_under_100ms']}",
            f"- glet VLA-adapted cases timed out after 30s: gr00t_two_model={ART['cases']['gr00t_two_model_gpulet']['timed_out']}, gr00t_four_model={ART['cases']['gr00t_four_model_gpulet']['timed_out']}, pi05_four_model={ART['cases']['pi05_four_model_scaledx10_gpulet']['timed_out']}",
        ],
        body_size=15,
    )
    slides.append(slide)

    # 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "为什么我们的方法更好")
    add_plain_rule(slide, Inches(1.2))
    add_lines(
        slide,
        Inches(0.8),
        Inches(1.45),
        Inches(11.8),
        Inches(5.6),
        [
            "- 我们不是只优化请求到来后的 compute scheduling，而是联合优化 future compute slots、resident model state、H2D/activation bandwidth、legal replan window。",
            "- Pi05 上，主要收益来自 25/50 control semantics + predictive residency/prefetch，admission 从 22.33 提升到 32.67，而 p95 仍约 43.21ms。",
            "- GR00T 上，主要收益来自 shared-prefix residency + phase-lock batching + fair admission：16 robots 时稳定在 58.43ms，reply-over 从 15 降到 0。",
            "- 这说明同一个 VLA-vGPU 抽象可以覆盖两类不同的 VLA family，只是资源重心不同。",
        ],
        size=20,
    )
    slides.append(slide)

    # 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "为什么这个方法可以 goes general")
    add_plain_rule(slide, Inches(1.2))
    add_section_block(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.8),
        "它不依赖具体模型实现",
        [
            "- 不依赖具体 transformer backbone",
            "- 不依赖具体 tokenizer 或 action head",
            "- 不依赖特定 checkpoint 格式",
            "- 不依赖特定厂商 runtime",
        ],
    )
    add_section_block(
        slide,
        Inches(6.75),
        Inches(1.45),
        Inches(5.8),
        "它只要求可测的工作负载属性",
        [
            "- chunk size",
            "- horizon process / legal replan window",
            "- single-request latency",
            "- same-model batch curve",
            "- resident state footprint",
            "- future request predictability",
            "- 然后选择 ResidentPrefetch / SharedPrefixPhaseBatch / Hybrid",
        ],
    )
    slides.append(slide)

    # 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "结论")
    add_plain_rule(slide, Inches(1.2))
    add_lines(
        slide,
        Inches(0.8),
        Inches(1.6),
        Inches(11.8),
        Inches(5.4),
        [
            "- VLA 需要的不是传统请求式 GPU 虚拟化，而是面向 Robot Lease 的时空联合 GPU 虚拟化。",
            "- Pi05 证明了：legal replan window + predictive residency/prefetch 可以稳定扩容，同时不增加端到端推理时延。",
            "- GR00T N1.6 证明了：shared-prefix residency + phase-lock batching + fair admission 是有效的。",
            "- 传统 GPU 虚拟化 / 资源划分方案不足的根本原因，是它们没有把模型状态、带宽和控制窗口当成一等资源。",
            "- 因此，VLA-vGPU 是一个既能解释 Pi05，又能解释 GR00T，并且可以 general 到 chunked-action VLA family 的系统抽象。",
        ],
        size=21,
    )
    slides.append(slide)

    total = len(slides)
    for idx, slide in enumerate(slides, start=1):
        add_page_num(slide, idx, total)

    prs.save(OUT)
    return OUT, total


def main():
    out, total = build_presentation()
    print(json.dumps({"out": str(out), "slides": total}, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
