#!/usr/bin/env python3

from __future__ import annotations

import json
import math
import tempfile
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.ticker import ScalarFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/root/autodl-tmp/VLAServing")
RESULTS = ROOT / "results"
DOCS = ROOT / "docs"
OUT = DOCS / "VLA_Workload_GPU_Virtualization_20260413.pptx"


COLORS = {
    "bg": RGBColor(247, 244, 238),
    "paper": RGBColor(255, 255, 255),
    "text": RGBColor(15, 23, 42),
    "muted": RGBColor(71, 85, 105),
    "teal": RGBColor(15, 118, 110),
    "teal_dark": RGBColor(11, 59, 54),
    "orange": RGBColor(194, 65, 12),
    "gold": RGBColor(180, 83, 9),
    "green": RGBColor(22, 163, 74),
    "red": RGBColor(220, 38, 38),
    "blue": RGBColor(37, 99, 235),
    "gray": RGBColor(203, 213, 225),
}


def load_json(name: str):
    with (RESULTS / name).open() as f:
        return json.load(f)


VS = load_json("unified_chunked_vla_vs_baselines_20260412.json")
EFF = load_json("unified_chunked_vla_effectiveness_20260412.json")
PI05 = load_json("pi05_four_model_residency_prefetch_system_20260406.json")
GR00T = load_json("gr00t_shared_prefix_phase_lock_batch_mps_20260412.json")
ART = load_json("public_gpu_serving_artifacts_20260411.json")


def rgb(name: str) -> RGBColor:
    return COLORS[name]


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb("bg")


def add_footer(slide, index: int, total: int):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"RTX 4090 24GB | Pi05 + GR00T N1.6 | VLA GPU Virtualization | {index}/{total}"
    p.font.size = Pt(10)
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = rgb("muted")
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title: str, subtitle: str | None = None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = rgb("teal")
    band.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(11.9), Inches(0.65))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = rgb("text")
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.12), Inches(11.7), Inches(0.5))
        stf = sub.text_frame
        stf.clear()
        p = stf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(13)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = rgb("muted")
        p.alignment = PP_ALIGN.LEFT


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: int = 18,
    color: str = "text",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = rgb(color)
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, bullets: list[str], size: int = 18, color: str = "text"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(8)
        p.bullet = True
    return box


def add_card(slide, left, top, width, height, title: str, body: list[str], accent: str = "teal"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("paper")
    shape.line.color.rgb = rgb("gray")
    shape.line.width = Pt(1.2)

    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.10))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = rgb(accent)
    accent_bar.line.fill.background()

    add_textbox(slide, left + Inches(0.18), top + Inches(0.18), width - Inches(0.3), Inches(0.35), title, 16, "text", True)
    add_bullets(slide, left + Inches(0.18), top + Inches(0.55), width - Inches(0.32), height - Inches(0.65), body, 13)


def add_metric_card(slide, left, top, width, height, label: str, value: str, accent: str = "teal"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("paper")
    shape.line.color.rgb = rgb("gray")
    shape.line.width = Pt(1.0)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.10))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(accent)
    band.line.fill.background()

    add_textbox(slide, left + Inches(0.16), top + Inches(0.16), width - Inches(0.2), Inches(0.25), label, 12, "muted", False)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.45), width - Inches(0.2), Inches(0.5), value, 22, accent, True)


def add_table(slide, left, top, width, height, headers: list[str], rows: list[list[str]], col_widths: list[float]):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    total = sum(col_widths)
    for idx, frac in enumerate(col_widths):
        table.columns[idx].width = int(width * (frac / total))
    header_fill = rgb("teal_dark")
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.name = "Microsoft YaHei"
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("paper") if i % 2 else RGBColor(250, 250, 250)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "Microsoft YaHei"
                p.font.color.rgb = rgb("text")
                p.alignment = PP_ALIGN.LEFT


def save_baseline_latency_chart(path: Path):
    pi05 = {
        "Ours": VS["pi05"]["latest_method"]["service_e2e_p95_ms"],
        "Full resident": VS["pi05"]["generic_full_resident_upper_bound"]["mean_latency_p95_ms"],
        "Clockwork": VS["pi05"]["best_conventional_baseline"]["mean_latency_p95_ms"],
        "REEF-like": VS["pi05"]["all_request_level_baselines"]["reef_like_temporal"]["mean_latency_p95_ms"],
        "DistServe": VS["pi05"]["all_request_level_baselines"]["distserve_like"]["mean_latency_p95_ms"],
    }
    gr00t = {
        "Ours": VS["gr00t_n1d6"]["latest_method_4_robot"]["mean_request_to_result_p95_ms"],
        "Full resident": VS["gr00t_n1d6"]["generic_full_resident_upper_bound"]["mean_latency_p95_ms"],
        "USHER-like": VS["gr00t_n1d6"]["best_conventional_baseline"]["mean_latency_p95_ms"],
        "REEF-like": VS["gr00t_n1d6"]["all_request_level_baselines"]["reef_like_temporal"]["mean_latency_p95_ms"],
        "DistServe": VS["gr00t_n1d6"]["all_request_level_baselines"]["distserve_like"]["mean_latency_p95_ms"],
    }

    fig, axes = plt.subplots(1, 2, figsize=(11.2, 3.7), dpi=200)
    for ax, title, data in zip(axes, ["Pi05", "GR00T N1.6"], [pi05, gr00t]):
        labels = list(data.keys())
        vals = list(data.values())
        bars = ax.bar(labels, vals, color=["#0f766e", "#0b3b36", "#2563eb", "#c2410c", "#7c2d12"])
        ax.set_yscale("log")
        ax.axhline(100, color="#dc2626", linestyle="--", linewidth=1.5)
        ax.set_title(title, fontsize=13, fontweight="bold")
        ax.set_ylabel("p95 latency (ms)")
        ax.set_ylim(30, 4000)
        ax.yaxis.set_major_formatter(ScalarFormatter())
        ax.tick_params(axis="x", labelrotation=22, labelsize=9)
        ax.tick_params(axis="y", labelsize=9)
        for bar, val in zip(bars, vals):
            ax.text(bar.get_x() + bar.get_width() / 2, val * 1.08, f"{val:.1f}", ha="center", va="bottom", fontsize=8)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="#f7f4ee")
    plt.close(fig)


def save_gr00t_scaling_chart(path: Path):
    scenarios = GR00T["scenarios"]
    xs = [4, 8, 16, 24]
    ys = [
        scenarios["1x_per_model"]["batch_only"]["best"]["metrics"]["mean_request_to_result_p95_ms"],
        EFF["gr00t_phase_lock_effectiveness"]["cases"][0]["phase_lock_batch"]["mean_request_to_result_p95_ms"],
        scenarios["4x_per_model"]["batch_only"]["best"]["metrics"]["mean_request_to_result_p95_ms"],
        EFF["gr00t_batch_mps_effectiveness"]["cases"][1]["batch_only"]["mean_request_to_result_p95_ms"],
    ]
    stable = [True, True, True, False]
    reply_over = [0, 0, 0, 18]
    fig, ax = plt.subplots(figsize=(6.0, 3.6), dpi=200)
    ax.plot(xs, ys, color="#0f766e", linewidth=2.5, marker="o", markersize=7)
    ax.axhline(100, color="#dc2626", linestyle="--", linewidth=1.5)
    ax.set_xlabel("Robots served simultaneously")
    ax.set_ylabel("p95 request-to-result (ms)")
    ax.set_title("GR00T scaling with shared-prefix + phase-lock")
    ax.set_xticks(xs)
    ax.set_ylim(35, 110)
    ax.grid(axis="y", alpha=0.25)
    for x, y, ok, ro in zip(xs, ys, stable, reply_over):
        label = "stable" if ok else f"unstable, over={ro}"
        color = "#16a34a" if ok else "#dc2626"
        ax.scatter([x], [y], color=color, s=45, zorder=3)
        ax.text(x, y + 3.0, f"{y:.1f}\n{label}", ha="center", va="bottom", fontsize=8)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="#f7f4ee")
    plt.close(fig)


def save_prefetch_chart(path: Path):
    fixed = EFF["pi05_prefetch_effectiveness"]["fixed_four_robot_level"]
    adm = EFF["pi05_prefetch_effectiveness"]["admission_level"]
    fig, axes = plt.subplots(1, 2, figsize=(8.0, 3.0), dpi=200)

    axes[0].bar(
        ["No prefetch", "Predictive prefetch"],
        [
            fixed["without_prefetch"]["mean_hard_miss_count"],
            fixed["with_prefetch"]["mean_hard_miss_count"],
        ],
        color=["#c2410c", "#0f766e"],
    )
    axes[0].set_title("Pi05 fixed 4 robots")
    axes[0].set_ylabel("hard misses")
    axes[0].set_ylim(0, 1.3)

    axes[1].bar(
        ["No prefetch", "Predictive prefetch"],
        [
            adm["without_prefetch"]["hard_miss_count"],
            adm["with_prefetch"]["hard_miss_count"],
        ],
        color=["#c2410c", "#0f766e"],
    )
    axes[1].set_title("Pi05 admission-level validation")
    axes[1].set_ylabel("hard misses")
    axes[1].set_ylim(0, 3.6)

    for ax in axes:
        ax.grid(axis="y", alpha=0.25)
        for p in ax.patches:
            ax.text(p.get_x() + p.get_width() / 2, p.get_height() + 0.05, f"{p.get_height():.0f}", ha="center", fontsize=9)

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="#f7f4ee")
    plt.close(fig)


def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total_slides = 11

    with tempfile.TemporaryDirectory() as tmpdir_str:
        tmpdir = Path(tmpdir_str)
        baseline_chart = tmpdir / "baseline_latency.png"
        gr00t_chart = tmpdir / "gr00t_scaling.png"
        prefetch_chart = tmpdir / "pi05_prefetch.png"
        save_baseline_latency_chart(baseline_chart)
        save_gr00t_scaling_chart(gr00t_chart)
        save_prefetch_chart(prefetch_chart)

        # Slide 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "VLA Workload 的 GPU 虚拟化", "面向 Pi05 与 GR00T N1.6 的单 GPU Serving 抽象、方法、实验与对照")
        add_textbox(slide, Inches(0.6), Inches(1.6), Inches(7.2), Inches(1.1),
                    "核心命题：VLA 不能只虚拟化 compute。必须联合虚拟化未来时间线、模型状态、带宽与合法重规划窗口。", 24, "text", True)
        add_card(slide, Inches(0.7), Inches(3.0), Inches(3.8), Inches(2.0),
                 "问题", ["多机器人绑定多微调模型", "100ms 内返回新 chunk", "AutoHorizon 使请求相位可预测但不固定"], "orange")
        add_card(slide, Inches(4.8), Inches(3.0), Inches(3.8), Inches(2.0),
                 "方法", ["Pi05: exact residency + predictive prefetch", "GR00T: shared-prefix + phase-lock batching", "admission/fairness 纳入统一控制面"], "teal")
        add_card(slide, Inches(8.9), Inches(3.0), Inches(3.7), Inches(2.0),
                 "结果", ["Pi05: 43.21ms p95, 0 hard miss", "GR00T: 4 机器人 43.88ms, 16 机器人 58.43ms", "传统 GPU 虚拟化 / 资源划分方法明显不足"], "blue")
        add_footer(slide, 1, total_slides)

        # Slide 2
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "为什么 VLA 不同于传统 Serving", "VLA workload 的瓶颈不是单一 queueing，而是闭环控制 + 模型状态 + 可预测性")
        add_bullets(slide, Inches(0.6), Inches(1.45), Inches(5.8), Inches(4.8), [
            "Next request predictable：当前 chunk 发出后，下一次请求大致何时发生是可预测的。",
            "Deadline is action exhaustion：真正的 deadline 不是平均响应时间，而是在旧 action chunk 耗尽前返回新 chunk。",
            "Model state is part of scheduling：多微调模型时，权重驻留、预取、swap、decode/apply 本身就是一等资源。",
            "Horizon is controllable：AutoHorizon/早重规划意味着系统可以在合法窗口内主动移动请求相位。"
        ], 18)
        add_card(slide, Inches(6.8), Inches(1.45), Inches(2.7), Inches(1.6), "传统抽象", ["queue", "space partition", "memory capacity", "SLO miss"], "orange")
        add_card(slide, Inches(9.8), Inches(1.45), Inches(2.8), Inches(1.6), "VLA 需要", ["predictive timeline", "resident model state", "copy/decode bandwidth", "legal replan window"], "teal")
        add_textbox(slide, Inches(6.9), Inches(3.5), Inches(5.2), Inches(1.1),
                    "结论：VLA 的虚拟化对象不是静态 CUDA context，也不是固定空间分区，而是一个面向闭环控制的可预测 lease。", 20, "teal_dark", True)
        add_footer(slide, 2, total_slides)

        # Slide 3
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "VLA-vGPU 抽象", "一个机器人占用的不是一份静态 GPU，而是一份时空联合 lease")
        add_textbox(slide, Inches(0.65), Inches(1.35), Inches(12.0), Inches(0.55),
                    "VLA-vGPU = <T, S, M, X, W, B, F>", 28, "teal_dark", True, PP_ALIGN.CENTER)
        cards = [
            ("T", "Temporal compute lease", ["未来 compute slots", "predictive reservation"], "teal"),
            ("S", "Spatial compute share", ["独占/MPS/batch lane", "必要时才分区"], "blue"),
            ("M", "Model-state residency", ["full shell", "shared prefix", "delta/hot pages"], "orange"),
            ("X", "Transfer/decode budget", ["H2D", "decode/apply", "activation buffer"], "gold"),
            ("W", "Legal request window", ["Pi05: [25,50]", "GR00T: [8,16]"], "green"),
            ("B", "Batch affinity", ["只和同模型聚合", "形成长期 cohort"], "teal"),
            ("F", "Fairness / admission", ["公平权重", "防止 batch bias"], "red"),
        ]
        x_positions = [0.6, 2.45, 4.3, 6.15, 8.0, 9.85, 11.2]
        widths = [1.7, 1.7, 1.7, 1.7, 1.7, 1.2, 1.4]
        for (tag, title, body, accent), x, w in zip(cards, x_positions, widths):
            add_card(slide, Inches(x), Inches(2.1), Inches(w), Inches(2.6), f"{tag} | {title}", body, accent)
        add_textbox(slide, Inches(0.8), Inches(5.2), Inches(12.0), Inches(0.9),
                    "抽象单位是 Robot Lease，而不是 Request：admit 的对象是一个长期闭环机器人，它绑定模型、频率、未来时间线与驻留价值。", 19, "text", True)
        add_footer(slide, 3, total_slides)

        # Slide 4
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "系统控制面与数据面", "统一控制面 + 模型家族特化数据面")
        # control plane boxes
        control_titles = ["Timeline Predictor", "Residency Manager", "Transfer Scheduler", "Compute Scheduler", "Admission Controller"]
        cx = [0.7, 3.15, 5.6, 8.05, 10.5]
        for title, x in zip(control_titles, cx):
            add_card(slide, Inches(x), Inches(1.65), Inches(2.2), Inches(1.45), title, [], "teal")
        for x in [2.78, 5.23, 7.68, 10.13]:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.1), Inches(0.28), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb("teal")
            arrow.line.fill.background()
        add_textbox(slide, Inches(0.8), Inches(3.45), Inches(12.0), Inches(0.35), "Data plane: compute stream + prefetch stream + decode/apply stream", 18, "teal_dark", True, PP_ALIGN.CENTER)
        # three stream lanes
        lanes = [
            ("compute stream", "当前 batch / 当前 block 推理", "teal"),
            ("prefetch stream", "后台 H2D / page fetch / next-use preparation", "blue"),
            ("decode/apply stream", "delta decode / page activation / state apply", "orange"),
        ]
        ly = [4.1, 5.0, 5.9]
        for (name, body, accent), y in zip(lanes, ly):
            add_card(slide, Inches(0.9), Inches(y), Inches(11.5), Inches(0.65), name, [body], accent)
        add_textbox(slide, Inches(0.95), Inches(6.75), Inches(11.8), Inches(0.25),
                    "Backend A: ExactDeltaPrefetch (Pi05) | Backend B: SharedPrefixPhaseBatch (GR00T)", 14, "muted", False, PP_ALIGN.CENTER)
        add_footer(slide, 4, total_slides)

        # Slide 5
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "Pi05: memory-state virtualization", "three-shell residency + predictive prefetch + H2D/decode overlap")
        add_card(slide, Inches(0.65), Inches(1.5), Inches(4.2), Inches(3.4), "Runtime design", [
            "Shell-A: 30Hz official FT 常驻",
            "Shell-B: 20Hz quantiles 常驻",
            "Shell-C: 两个 10Hz 模型共享",
            "exact shared-base / delta pages 减少显存占用",
            "利用 AutoHorizon 预测 next-use，提前 prefetch"
        ], "teal")
        # Shell diagram
        for i, (label, col) in enumerate([("Shell-A", "teal"), ("Shell-B", "blue"), ("Shell-C", "orange")]):
            box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.25 + i * 2.15), Inches(1.85), Inches(1.75), Inches(1.55))
            box.fill.solid()
            box.fill.fore_color.rgb = rgb(col)
            box.line.fill.background()
            add_textbox(slide, Inches(5.4 + i * 2.15), Inches(2.1), Inches(1.45), Inches(0.32), label, 16, "paper", True, PP_ALIGN.CENTER)
            desc = {
                "Shell-A": "30Hz\nresident",
                "Shell-B": "20Hz\nresident",
                "Shell-C": "10Hz / 10Hz\nshared shell",
            }[label]
            add_textbox(slide, Inches(5.37 + i * 2.15), Inches(2.5), Inches(1.48), Inches(0.55), desc, 13, "paper", False, PP_ALIGN.CENTER)
        add_metric_card(slide, Inches(5.25), Inches(4.15), Inches(2.15), Inches(1.1), "Measured infer time", "≈ 43.2 ms", "teal")
        add_metric_card(slide, Inches(7.55), Inches(4.15), Inches(2.15), Inches(1.1), "Full swap cost", "289.47 ms", "orange")
        add_metric_card(slide, Inches(9.85), Inches(4.15), Inches(2.15), Inches(1.1), "GPU memory", "22.455 GB", "blue")
        add_textbox(slide, Inches(5.25), Inches(5.55), Inches(6.8), Inches(0.85),
                    "关键点：如果不做预测式 residency / prefetch，289ms 的模型状态移动会直接破坏实时性；真正的收益来自把它隐藏到未来时间线里。", 16, "text", True)
        add_footer(slide, 5, total_slides)

        # Slide 6
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "Pi05 实验：预取有效，且明显优于旧方法", "4 个不同微调模型分别服务 30/20/10/10Hz 机器人")
        slide.shapes.add_picture(str(prefetch_chart), Inches(0.65), Inches(1.45), width=Inches(4.4))
        add_metric_card(slide, Inches(5.35), Inches(1.55), Inches(2.1), Inches(1.0), "p95 latency", "43.21 ms", "teal")
        add_metric_card(slide, Inches(7.65), Inches(1.55), Inches(2.1), Inches(1.0), "Hard misses", "0", "green")
        add_metric_card(slide, Inches(9.95), Inches(1.55), Inches(2.1), Inches(1.0), "Strict deadlines", "45–100 ms 全通过", "blue")
        add_card(slide, Inches(5.35), Inches(2.9), Inches(6.7), Inches(1.6), "Proof of effectiveness", [
            "固定 4 机器人：predictive prefetch 将 hard miss 1 -> 0，SLA miss 1 -> 0。",
            "admission-level：hard miss 3 -> 0，p95 43.205 -> 43.198 ms，说明收益来自正确的状态时序而不是额外 compute。"
        ], "teal")
        add_textbox(slide, Inches(5.35), Inches(4.9), Inches(6.8), Inches(0.65),
                    "对照：generic full-resident 104.92 ms；最佳传统 baseline Clockwork-like 364.84 ms。", 18, "orange", True)
        add_footer(slide, 6, total_slides)

        # Slide 7
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "GR00T N1.6: temporal cohort virtualization", "shared-prefix residency + phase-lock batching + quota-fair admission")
        add_card(slide, Inches(0.65), Inches(1.55), Inches(4.3), Inches(3.2), "Runtime design", [
            "所有微调模型通过 shared prefix 常驻",
            "只允许 same-model batch，利用 batch latency curve",
            "在 [8,16] action 合法窗口内做 early replan / phase correction",
            "用 quota-fair admission 避免 greedy batch-first 偏置"
        ], "teal")
        add_textbox(slide, Inches(5.25), Inches(1.75), Inches(6.7), Inches(0.45),
                    "相位控制目标：把未来请求移动成长期稳定 cohort，而不是只优化当前一次 queue。", 18, "teal_dark", True)
        # mini cohort diagram
        for row, model in enumerate(["30Hz bridge", "20Hz fractal", "10Hz libero", "10Hz rel30k"]):
            add_textbox(slide, Inches(5.35), Inches(2.3 + row * 0.65), Inches(1.8), Inches(0.25), model, 13, "muted", True)
            for col in range(4):
                box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0 + col * 1.0), Inches(2.18 + row * 0.65), Inches(0.78), Inches(0.34))
                box.fill.solid()
                box.fill.fore_color.rgb = rgb("teal") if row < 2 else (rgb("orange") if row == 2 else rgb("blue"))
                box.line.fill.background()
        add_textbox(slide, Inches(7.0), Inches(5.2), Inches(4.6), Inches(0.55),
                    "固定 cohort 后，同模型请求自然汇聚成 batch 1 / 2 / 4 / 6。", 16, "text", True)
        add_metric_card(slide, Inches(5.35), Inches(5.85), Inches(2.0), Inches(0.95), "4 robots", "43.88 ms", "teal")
        add_metric_card(slide, Inches(7.55), Inches(5.85), Inches(2.0), Inches(0.95), "16 robots", "58.43 ms", "blue")
        add_metric_card(slide, Inches(9.75), Inches(5.85), Inches(2.0), Inches(0.95), "16-robot batch", "4.0", "orange")
        add_footer(slide, 7, total_slides)

        # Slide 8
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "GR00T 实验：phase-lock 与 fairness 的有效性", "这不是泛化技巧，而是直接决定能否稳定 serving")
        slide.shapes.add_picture(str(gr00t_chart), Inches(0.65), Inches(1.45), width=Inches(5.7))
        add_card(slide, Inches(6.7), Inches(1.55), Inches(5.9), Inches(1.45), "Phase-lock proof", [
            "8 robots: strict horizon reply-over 7 -> 0，batch 1.38 -> 2.00。",
            "16 robots: strict horizon reply-over 15 -> 0，batch 2.20 -> 4.00。"
        ], "teal")
        add_card(slide, Inches(6.7), Inches(3.2), Inches(5.9), Inches(1.35), "Fair admission proof", [
            "accept-rate gap: 0.2921 -> 0.0994",
            "final-count gap: 18 -> 2",
            "说明 batch-aware admission 若不加公平约束，会系统性偏置某些模型/频率。"
        ], "orange")
        add_card(slide, Inches(6.7), Inches(4.8), Inches(5.9), Inches(1.25), "MPS ablation", [
            "16 robots: batch-only 与 batch+MPS 完全相同。",
            "24 robots: 两者都不稳定，reply-over 都是 18。"
        ], "blue")
        add_footer(slide, 8, total_slides)

        # Slide 9
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "为什么之前的 GPU 虚拟化 / 资源划分方法不够好", "不仅假设不匹配，而且实验上确实做不好 VLA")
        headers = ["方法", "核心假设", "VLA mismatch", "Pi05 / GR00T 结果"]
        rows = [
            ["GPUlet-like", "静态 partition + offline schedule", "不显式管理模型状态与合法重规划窗口", "三种 partition 在 100ms 下均 infeasible"],
            ["REEF-like", "temporal switching", "忽略 future model-state lease", "752.7 ms / 1195.2 ms p95"],
            ["Clockwork-like", "固定 clock + shell pool", "无法利用可预测相位与 AutoHorizon", "364.8 ms / 437.0 ms p95"],
            ["USHER-like", "空间切分优先", "没有预测式 residency/prefetch", "Pi05 471.7 ms；GR00T 520.7 ms"],
            ["DistServe-like", "prefill/decode 解耦", "VLA 不是自然的 prefill/decode workload", "1980.1 ms / 2689.9 ms p95"],
            ["Oracle full resident", "所有模型都在显存中", "即便不换权重，时间线冲突仍存在", "104.9 ms / 140.3 ms p95"],
        ]
        add_table(slide, Inches(0.55), Inches(1.55), Inches(12.15), Inches(4.95), headers, rows, [1.4, 2.1, 3.0, 2.2])
        add_textbox(slide, Inches(0.75), Inches(6.75), Inches(11.8), Inches(0.3),
                    "结论：VLA 不能只做空间划分，也不能只做时间片；必须把时间、模型状态、带宽、合法窗口联合起来虚拟化。", 15, "teal_dark", True)
        add_footer(slide, 9, total_slides)

        # Slide 10
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "公开代码 / Artifact 角度的补充证据", "不仅 paper-faithful reproduction 不适配，公开代码路径也直接暴露出 VLA mismatch")
        add_card(slide, Inches(0.65), Inches(1.55), Inches(4.0), Inches(2.6), "GPUlet / glet public code", [
            "sanity_single_vgg16 可运行：exit_code=0，说明二进制本身没问题。",
            "但 gr00t_two_model / gr00t_four_model / pi05_four_model_scaledx10 三个 VLA-adapted case 全部 30s timeout，且无输出文件。"
        ], "teal")
        add_card(slide, Inches(4.9), Inches(1.55), Inches(3.7), Inches(2.6), "Pi05 request-rate mismatch", [
            "Pi05 mean request RPS = [1.02, 0.68, 0.34, 0.34]",
            "glet 只接受整数 RPS",
            "最近整数近似会退化成 [1, 1, 0, 0]，直接丢掉两个 10Hz 机器人"
        ], "orange")
        add_card(slide, Inches(8.9), Inches(1.55), Inches(3.8), Inches(2.6), "REEF / Paella artifact gap", [
            "REEF 依赖 MI50 + ROCm 4.3 + 定制 amdgpu/hip",
            "Paella 依赖 TVM/LLIS 自定义编译栈",
            "都不是当前 Pi05 / GR00T VLA 栈的直接运行路径"
        ], "blue")
        add_textbox(slide, Inches(0.75), Inches(4.65), Inches(11.6), Inches(1.0),
                    "因此“之前的方法不够好”不只是概念判断，也不是调参问题：它们的资源模型与输入模型就没有把 VLA 的模型状态、相位可预测性和 action-exhaustion deadline 放进系统抽象。", 18, "text", True)
        add_footer(slide, 10, total_slides)

        # Slide 11
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_title(slide, "总结", "VLA Workload 专用 GPU 虚拟化的最终观点")
        add_bullets(slide, Inches(0.75), Inches(1.5), Inches(8.3), Inches(4.8), [
            "VLA-vGPU 的核心不是静态 partition，而是 Robot Lease：<T, S, M, X, W, B, F>。",
            "Pi05 的关键在 memory-state virtualization：three-shell、exact/shared-base、predictive prefetch、bandwidth scheduling。",
            "GR00T 的关键在 temporal cohort virtualization：shared-prefix、phase-lock batching、quota-fair admission。",
            "实验上，Pi05 达到 43.21ms p95 / 0 hard miss；GR00T 4 机器人 43.88ms，16 机器人 58.43ms。",
            "传统 GPU 虚拟化 / 资源划分方法没有把 future timeline、resident model state、copy/decode bandwidth、legal replan window 联合建模，因此不适合 VLA。"
        ], 19)
        add_card(slide, Inches(9.35), Inches(1.8), Inches(3.0), Inches(1.6), "一句话", [
            "VLA 专用 GPU 虚拟化 = 未来 compute 时间片 + 显存驻留份额 + 状态搬运带宽 + 合法重规划窗口 + batch 亲和性。"
        ], "teal")
        add_card(slide, Inches(9.35), Inches(3.8), Inches(3.0), Inches(1.6), "可直接引用", [
            "docs/VLA_WORKLOAD_GPU_VGPU_ABSTRACTION.md",
            "results/unified_chunked_vla_vs_baselines_20260412.json",
            "results/unified_chunked_vla_effectiveness_20260412.json",
        ], "orange")
        add_footer(slide, 11, total_slides)

        prs.save(OUT)


if __name__ == "__main__":
    create_deck()
    print(OUT)
